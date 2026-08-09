"""Document text extractor — modernised 2026-05-26.

Supports PDF, DOCX, PPTX, XLSX, CSV, TXT, MD. Pure-Python libraries
(Lambda-safe, no native deps). Robust against non-UTF-8 encodings via
chardet auto-detection.

Library choices (all current and actively maintained as of 2026):
    PDF     → pypdf       (drop-in successor to retired PyPDF2)
    DOCX    → python-docx (paragraphs + tables)
    PPTX    → python-pptx (slide shapes + speaker notes)
    XLSX    → openpyxl    (read-only mode for speed)
    CSV     → stdlib csv  (auto-detected encoding)
    TXT/MD  → stdlib decode (auto-detected encoding)

Public surface (unchanged for backward compat):
    extract_text_from_file(file_bytes: bytes, filename: str) -> str
    process_multiple_files(files_data: list[tuple[str, bytes]]) -> str
"""

from __future__ import annotations

import csv
import io
import logging
from typing import Iterable, List, Tuple

logger = logging.getLogger("pipelyt.docs")


# ─── Extension dispatch ─────────────────────────────────────────────────────


def _ext(filename: str) -> str:
    if "." not in filename:
        return ""
    return filename.rsplit(".", 1)[-1].lower().strip()


# ─── Encoding detection ─────────────────────────────────────────────────────


def _decode_bytes(content: bytes) -> str:
    """Try UTF-8 first, then chardet auto-detect, then fall back to latin-1
    (which always succeeds — every byte sequence is a valid latin-1 string).
    Logs a warning if the detected encoding isn't UTF-8 so operators can spot
    files that may have transliteration issues.
    """
    if not content:
        return ""
    # Fast path
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        pass
    # Auto-detect
    try:
        import chardet  # lazy import — only pulled when needed
        detected = chardet.detect(content[:65536])  # 64KB sample is enough
        enc = (detected or {}).get("encoding") or ""
        conf = (detected or {}).get("confidence") or 0.0
        if enc and conf > 0.5:
            logger.info(
                "Detected non-UTF-8 file encoding: %s (confidence %.2f)", enc, conf,
            )
            try:
                return content.decode(enc, errors="replace")
            except (UnicodeDecodeError, LookupError):
                pass
    except ImportError:
        logger.debug("chardet is not installed; falling back to latin-1 for non-UTF-8 bytes")
    # Last resort — never fails, may produce odd chars for non-Latin scripts.
    return content.decode("latin-1", errors="replace")


# ─── Per-format extractors ──────────────────────────────────────────────────


def _extract_pdf(content: bytes) -> str:
    """Extract text from a PDF using pypdf (replaces deprecated PyPDF2).

    Iterates pages, calls page.extract_text() on each. Empty/None per-page
    returns are skipped silently — common for image-only PDFs that would
    otherwise need OCR (out of scope here).
    """
    from pypdf import PdfReader  # lazy import
    reader = PdfReader(io.BytesIO(content))
    parts: List[str] = []
    for i, page in enumerate(reader.pages):
        try:
            t = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Could not extract text from PDF page %d: %s "
                "(continuing with the remaining pages)",
                i, e,
            )
            t = ""
        if t.strip():
            parts.append(t)
    return "\n\n".join(parts)


def _extract_docx(content: bytes) -> str:
    """Extract text from a .docx file (Word).

    Reads both paragraphs and tables (tables were silently dropped by the old
    code, which lost meaningful structured content in spec docs / pricing
    sheets exported to Word).
    """
    from docx import Document  # python-docx — lazy import
    doc = Document(io.BytesIO(content))
    parts: List[str] = []
    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Tables — each row joined by tab, each cell by its text
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    """Extract text from a .pptx deck (PowerPoint).

    Pulls text from every shape on every slide PLUS speaker notes. Slides
    are separated by `--- Slide N ---` markers so the chunker preserves
    slide structure when paragraph-splitting.
    """
    from pptx import Presentation  # python-pptx — lazy import
    prs = Presentation(io.BytesIO(content))
    parts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        # Visible text on the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            # Tables nested inside slides
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    """Extract text from a .xlsx workbook (Excel).

    Iterates every sheet, every row. Empty rows/cells are skipped. Read-only
    mode keeps memory low even for big workbooks. Sheet boundaries marked
    with `--- Sheet: <name> ---` so the chunker keeps tabs separated.
    """
    from openpyxl import load_workbook  # lazy import
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: List[str] = []
    for sheet in wb.worksheets:
        parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_csv(content: bytes) -> str:
    """Extract text from a .csv file. Detects encoding first (handles non-UTF-8
    Excel exports), then uses stdlib csv module so we don't pull pandas just
    for this."""
    text = _decode_bytes(content)
    reader = csv.reader(io.StringIO(text))
    rows: List[str] = []
    for row in reader:
        cleaned = [cell.strip() for cell in row if cell and cell.strip()]
        if cleaned:
            rows.append(", ".join(cleaned))
    return "\n".join(rows)


def _extract_text(content: bytes) -> str:
    """Plain text / markdown extractor — encoding-tolerant."""
    return _decode_bytes(content)


# ─── Public API ─────────────────────────────────────────────────────────────


_EXTRACTORS = {
    "pdf":  _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "ppt":  _extract_pptx,   # legacy .ppt routes through pptx best-effort
    "xlsx": _extract_xlsx,
    "xls":  _extract_xlsx,   # legacy .xls routes through openpyxl best-effort
    "csv":  _extract_csv,
    "txt":  _extract_text,
    "md":   _extract_text,
}


def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """Extract plain text from a supported document format.

    Args:
        file_content: raw bytes of the file
        filename: original filename (used to dispatch on extension)

    Returns:
        Extracted plain text with normalised whitespace. Empty string if
        nothing extractable. Never raises — extraction errors are logged
        and an empty string is returned so the caller can decide whether
        to reject the upload.
    """
    ext = _ext(filename)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        # Unknown extension — try plain-text decode as last resort
        logger.info(
            "Unknown file extension '.%s' for '%s' — trying plain-text decode",
            ext, filename,
        )
        try:
            return _decode_bytes(file_content).strip()
        except Exception as exc:
            logger.error(
                "Plain-text fallback failed for '%s': %s", filename, exc,
            )
            return ""

    try:
        text = extractor(file_content)
    except Exception as exc:
        logger.exception(
            "Text extraction failed for '%s' (.%s): %s — returning empty string",
            filename, ext, exc,
        )
        return ""

    return (text or "").strip()


def process_multiple_files(files_data: Iterable[Tuple[str, bytes]]) -> str:
    """Concatenate extracted text from multiple files with file-boundary markers.

    Preserves the legacy signature used elsewhere in the codebase. Each file's
    content is wrapped in `--- Start of File: ... ---` / `--- End of File: ... ---`
    so downstream chunkers can preserve file boundaries when splitting.
    """
    combined: List[str] = []
    for filename, content in files_data:
        text = extract_text_from_file(content, filename)
        if not text:
            continue
        combined.append(f"--- Start of File: {filename} ---")
        combined.append(text)
        combined.append(f"--- End of File: {filename} ---")
    return "\n".join(combined)
