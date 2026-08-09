"""POST /nexus/kb/extract — multi-file text extraction for the
NewCampaign wizard's KB step.

Accepts a multipart upload of one or more .pdf / .docx / .pptx files.
For each file returns:
  - filename
  - char_count (length of extracted text)
  - error (None on success, short string on failure)

And a single combined `combined_text` blob the frontend can append into
the knowledge_base textarea.

PDF + DOCX are routed through apps/backend/services/doc_processor —
already battle-tested in the Pipelyt main app for AI-content generation.
PPTX is handled inline here (python-pptx).

Hard 25 MB per file ceiling. No persistence — this endpoint is a pure
extraction utility; chunking + embedding still happens via
/nexus/analyze when the user clicks Launch (knowledge_base field).
"""
from __future__ import annotations

import io
import logging
from typing import Any, Dict, List

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from sqlalchemy.orm import Session

from core.database import get_db
from models import User
from nexus.deps import (  # type: ignore[import-not-found]
    get_nexus_user,
    require_current_workspace,
)
from services.doc_processor import extract_text_from_file

logger = logging.getLogger("pipelyt.nexus.kb_extract")

router = APIRouter(prefix="/nexus/kb", tags=["Nexus — Knowledge Base"])

MAX_BYTES_PER_FILE = 25 * 1024 * 1024  # 25 MB
ALLOWED_EXTS = {".pdf", ".docx", ".pptx"}


def _extract_pptx(file_bytes: bytes) -> str:
    """Lazy-import python-pptx so missing dep gives a clear error.

    Walks every slide, every shape, and concatenates any text frames
    (including grouped shapes and table cells). Best-effort — silently
    skips shapes that don't have text.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:  # noqa: BLE001
        raise RuntimeError(
            "python-pptx is not installed. Run `pip install -r requirements-nexus.txt`."
        ) from exc

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_lines: List[str] = []
        for shape in slide.shapes:
            # Plain text shapes
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_lines.append(line)
            # Tables
            if getattr(shape, "has_table", False) and shape.table is not None:
                for row in shape.table.rows:
                    row_cells = [
                        (cell.text or "").strip() for cell in row.cells
                    ]
                    row_cells = [c for c in row_cells if c]
                    if row_cells:
                        slide_lines.append(" | ".join(row_cells))
            # Grouped shapes
            if shape.shape_type == 6 and hasattr(shape, "shapes"):  # 6 = MSO_SHAPE_TYPE.GROUP
                for sub in shape.shapes:
                    if getattr(sub, "has_text_frame", False) and sub.text_frame is not None:
                        for para in sub.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                slide_lines.append(line)
        if slide_lines:
            chunks.append(f"# Slide {slide_idx}")
            chunks.extend(slide_lines)
            chunks.append("")
    return "\n".join(chunks).strip()


def _extract_one(content: bytes, filename: str) -> str:
    """Dispatch by extension. Raises on unsupported types."""
    lower = (filename or "").lower()
    if lower.endswith(".pptx"):
        return _extract_pptx(content)
    # PDF/DOCX/CSV/plain — reuse the existing Pipelyt doc processor.
    return extract_text_from_file(content, filename)


@router.post("/extract")
async def kb_extract(
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db),
    user: User = Depends(get_nexus_user),
    workspace=Depends(require_current_workspace),
) -> Dict[str, Any]:
    """Accept N files, return per-file metadata + combined text blob."""
    if not files:
        raise HTTPException(status_code=400, detail="No files supplied.")

    per_file: List[Dict[str, Any]] = []
    combined_parts: List[str] = []

    for f in files:
        name = f.filename or "unnamed"
        ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""

        if ext not in ALLOWED_EXTS:
            per_file.append(
                {
                    "filename": name,
                    "char_count": 0,
                    "error": f"Unsupported file type ({ext or 'no extension'}). "
                    "Accepted: PDF, DOCX, PPTX.",
                }
            )
            continue

        try:
            body = await f.read()
        except Exception as exc:  # noqa: BLE001
            per_file.append(
                {"filename": name, "char_count": 0, "error": f"Read failed: {exc}"}
            )
            continue

        if len(body) > MAX_BYTES_PER_FILE:
            per_file.append(
                {
                    "filename": name,
                    "char_count": 0,
                    "error": f"File too large ({len(body)//1024//1024} MB). "
                    f"Max {MAX_BYTES_PER_FILE//1024//1024} MB per file.",
                }
            )
            continue

        try:
            text = _extract_one(body, name) or ""
        except Exception as exc:  # noqa: BLE001
            logger.exception("KB extract failed for %s: %s", name, exc)
            per_file.append(
                {"filename": name, "char_count": 0, "error": f"Extract failed: {exc}"}
            )
            continue

        per_file.append(
            {"filename": name, "char_count": len(text), "error": None}
        )
        if text:
            combined_parts.append(f"--- {name} ---\n{text}")

    return {
        "files": per_file,
        "combined_text": "\n\n".join(combined_parts),
    }
